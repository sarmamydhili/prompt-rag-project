"""PowerPoint presenter notes → OpenAI TTS MP3 files (Phase 1)."""

import os
import re
import sys
from pathlib import Path
from typing import Optional

from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation

# ---------------------------------------------------------------------------
# Settings — edit these before running
# ---------------------------------------------------------------------------

TTS_MODEL = "gpt-4o-mini-tts"
# Male: onyx, echo, fable, ash  |  Female/neutral: coral, nova, shimmer, sage
VOICE = "coral"
AUDIO_FORMAT = "mp3"
SPEECH_INSTRUCTIONS = (
    "Speak clearly and naturally in a friendly instructional tone. "
    "Use a comfortable teaching pace. Pause naturally between ideas and "
    "slow down slightly for important concepts. Avoid sounding like you "
    "are simply reading text."
)

REGENERATE_AUDIO = True
TEST_MODE = False
TEST_SLIDES = [1, 5, 10]

# OpenAI TTS input limit (characters)
MAX_TTS_INPUT_CHARS = 4096

# PowerPoint default placeholder phrases (case-insensitive match)
PLACEHOLDER_PATTERNS = [
    re.compile(r"^click to add notes\.?$", re.IGNORECASE),
    re.compile(r"^click to edit master text styles\.?$", re.IGNORECASE),
    re.compile(r"^click to edit notes\.?$", re.IGNORECASE),
]

NARRATION_LABEL_LINE = re.compile(r"^NARRATION:?\s*$", re.IGNORECASE)
NARRATION_LABEL_PREFIX = re.compile(r"^NARRATION:\s*", re.IGNORECASE)

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

BASE_DIR = Path(__file__).resolve().parent
INPUT_DIR = BASE_DIR / "input"
AUDIO_DIR = BASE_DIR / "audio"


def find_powerpoint() -> Path:
    """Find the single .pptx file in input/. Exit on 0 or >1 matches."""
    pptx_files = sorted(
        p for p in INPUT_DIR.glob("*.pptx") if not p.name.startswith("~$")
    )

    if not pptx_files:
        print(f"Error: No PowerPoint file found in {INPUT_DIR}")
        print("Place one .pptx file in the input/ folder and run again.")
        sys.exit(1)

    if len(pptx_files) > 1:
        print(f"Error: Found {len(pptx_files)} PowerPoint files in {INPUT_DIR}:")
        for f in pptx_files:
            print(f"  - {f.name}")
        print("Leave only one .pptx file in the input/ folder and run again.")
        sys.exit(1)

    return pptx_files[0]


def strip_narration_label(text: str) -> str:
    """Remove leading NARRATION: label from presenter notes."""
    if not text:
        return ""

    lines = text.splitlines()
    if not lines:
        return text

    if NARRATION_LABEL_LINE.match(lines[0].strip()):
        lines = lines[1:]
    elif NARRATION_LABEL_PREFIX.match(lines[0]):
        lines[0] = NARRATION_LABEL_PREFIX.sub("", lines[0], count=1)

    return "\n".join(lines)


def clean_narration(text: str) -> str:
    """Lightweight cleanup: trim, collapse whitespace, normalize blank lines."""
    if not text:
        return ""

    text = strip_narration_label(text)

    lines = [line.strip() for line in text.splitlines()]
    cleaned_lines = []
    prev_blank = False

    for line in lines:
        if not line:
            if not prev_blank:
                cleaned_lines.append("")
            prev_blank = True
            continue

        line = re.sub(r"[ \t]+", " ", line)
        cleaned_lines.append(line)
        prev_blank = False

    result = "\n".join(cleaned_lines).strip()
    return result


def is_placeholder_notes(text: str) -> bool:
    """Return True if text matches known PowerPoint placeholder content."""
    normalized = text.strip()
    if not normalized:
        return True

    for pattern in PLACEHOLDER_PATTERNS:
        if pattern.match(normalized):
            return True

    return False


def extract_slide_notes(presentation: Presentation):
    """Yield (slide_number, cleaned_text_or_None) for each slide in order."""
    for slide_number, slide in enumerate(presentation.slides, start=1):
        raw = ""

        if slide.has_notes_slide:
            text_frame = slide.notes_slide.notes_text_frame
            if text_frame is not None:
                raw = text_frame.text or ""

        cleaned = clean_narration(raw)

        if not cleaned or is_placeholder_notes(cleaned):
            yield slide_number, None
        else:
            yield slide_number, cleaned


def audio_path_for_slide(slide_number: int) -> Path:
    """Return output path: audio/slide_001.mp3"""
    return AUDIO_DIR / f"slide_{slide_number:03d}.{AUDIO_FORMAT}"


def generate_audio(client: OpenAI, text: str, output_path: Path) -> None:
    """Call OpenAI TTS and write the audio file."""
    with client.audio.speech.with_streaming_response.create(
        model=TTS_MODEL,
        voice=VOICE,
        input=text,
        instructions=SPEECH_INSTRUCTIONS,
        response_format=AUDIO_FORMAT,
    ) as response:
        response.stream_to_file(output_path)


def process_slide(
    client: OpenAI,
    slide_number: int,
    narration: Optional[str],
    stats: dict,
) -> None:
    """Process one slide: skip, reuse existing audio, or generate new audio."""
    stats["slides_processed"] += 1

    if narration is None:
        print(f"Slide {slide_number}: no narration - skipped")
        stats["slides_without_notes"] += 1
        return

    print(f"Slide {slide_number}: narration found")

    if len(narration) > MAX_TTS_INPUT_CHARS:
        print(
            f"Slide {slide_number}: error - narration exceeds "
            f"{MAX_TTS_INPUT_CHARS} character limit ({len(narration)} chars)"
        )
        stats["errors"] += 1
        return

    output_path = audio_path_for_slide(slide_number)

    if output_path.exists() and not REGENERATE_AUDIO:
        print(f"Slide {slide_number}: audio already exists - skipped")
        stats["existing_audio_skipped"] += 1
        return

    try:
        generate_audio(client, narration, output_path)
        print(f"Slide {slide_number}: audio generated -> {output_path.name}")
        stats["audio_generated"] += 1
    except Exception as exc:
        print(f"Slide {slide_number}: error - {exc}")
        stats["errors"] += 1


def main() -> None:
    load_dotenv(BASE_DIR / ".env")

    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key or api_key == "your-api-key-here":
        print("Error: OPENAI_API_KEY is not set in .env")
        sys.exit(1)

    INPUT_DIR.mkdir(exist_ok=True)
    AUDIO_DIR.mkdir(exist_ok=True)

    pptx_path = find_powerpoint()
    presentation = Presentation(str(pptx_path))
    total_slides_in_deck = len(presentation.slides)

    if TEST_MODE:
        target_slides = set(TEST_SLIDES)
        print(f"Test mode: processing slides {sorted(target_slides)}")
    else:
        target_slides = None
        print("Full mode: processing all slides")

    client = OpenAI()

    stats = {
        "slides_processed": 0,
        "audio_generated": 0,
        "existing_audio_skipped": 0,
        "slides_without_notes": 0,
        "errors": 0,
    }

    for slide_number, narration in extract_slide_notes(presentation):
        if target_slides is not None and slide_number not in target_slides:
            continue
        process_slide(client, slide_number, narration, stats)

    print()
    print(f"PowerPoint: {pptx_path.name}")
    print(f"Total slides in deck: {total_slides_in_deck}")

    if TEST_MODE:
        valid_test_slides = [s for s in TEST_SLIDES if 1 <= s <= total_slides_in_deck]
        print(f"Test slides requested: {TEST_SLIDES}")
        print(f"Test slides in range: {valid_test_slides}")

    print(f"Slides processed: {stats['slides_processed']}")
    print(f"Audio files generated: {stats['audio_generated']}")
    print(f"Existing audio skipped: {stats['existing_audio_skipped']}")
    print(f"Slides without notes: {stats['slides_without_notes']}")
    print(f"Errors: {stats['errors']}")


if __name__ == "__main__":
    main()
