"""Set timed slide advance based on MP3 narration length (Phase 3).

Reads the Phase 2 output from output/, sets each slide to auto-advance after
its MP3 duration (plus a small buffer). Keeps click-to-advance enabled by default.
Writes a new file — does not modify the Phase 2 output or input/ original.
"""

import re
import sys
from pathlib import Path
from typing import Dict

from lxml import etree
from mutagen.mp3 import MP3
from pptx import Presentation

# ---------------------------------------------------------------------------
# Settings — edit these before running
# ---------------------------------------------------------------------------

BUFFER_MS = 500
DEFAULT_SLIDE_DURATION_MS = 3000
ALLOW_CLICK_ADVANCE = True

TEST_MODE = False
TEST_SLIDES = [1, 5, 10]

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

BASE_DIR = Path(__file__).resolve().parent
AUDIO_DIR = BASE_DIR / "audio"
OUTPUT_DIR = BASE_DIR / "output"

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
NSMAP = {"p": P_NS}

AUDIO_FILE_PATTERN = re.compile(r"^slide_(\d+)\.mp3$", re.IGNORECASE)


def find_embedded_powerpoint() -> Path:
    """Find the single Phase 2 *_with_audio.pptx in output/."""
    pptx_files = sorted(
        p
        for p in OUTPUT_DIR.glob("*_with_audio.pptx")
        if not p.name.startswith("~$") and not p.stem.endswith("_with_audio_timed")
    )

    if not pptx_files:
        print(f"Error: No Phase 2 PowerPoint found in {OUTPUT_DIR}")
        print("Expected a file like MyDeck_with_audio.pptx. Run embed_audio.py first.")
        sys.exit(1)

    if len(pptx_files) > 1:
        print(f"Error: Found {len(pptx_files)} Phase 2 files in {OUTPUT_DIR}:")
        for f in pptx_files:
            print(f"  - {f.name}")
        print("Leave only one *_with_audio.pptx file and run again.")
        sys.exit(1)

    return pptx_files[0]


def find_audio_files() -> Dict[int, Path]:
    """Map slide numbers to MP3 paths from audio/slide_XXX.mp3."""
    audio_files: Dict[int, Path] = {}

    for path in sorted(AUDIO_DIR.glob("slide_*.mp3")):
        match = AUDIO_FILE_PATTERN.match(path.name)
        if match:
            audio_files[int(match.group(1))] = path

    return audio_files


def xpath(element, query: str):
    return etree.ElementBase.xpath(element, query, namespaces=NSMAP)


def mp3_duration_ms(audio_path: Path) -> int:
    """Return MP3 duration in milliseconds."""
    audio = MP3(str(audio_path))
    return int(audio.info.length * 1000)


def advance_duration_ms(slide_number: int, audio_files: Dict[int, Path]) -> int:
    """Compute auto-advance time for a slide."""
    audio_path = audio_files.get(slide_number)
    if audio_path is None:
        return DEFAULT_SLIDE_DURATION_MS

    return mp3_duration_ms(audio_path) + BUFFER_MS


def set_slide_advance(slide, duration_ms: int) -> None:
    """Set or update the slide transition to auto-advance after duration_ms."""
    slide_element = slide._element
    adv_click = "1" if ALLOW_CLICK_ADVANCE else "0"

    transitions = xpath(slide_element, ".//p:transition")
    if transitions:
        transition = transitions[0]
    else:
        clr_map_nodes = xpath(slide_element, ".//p:clrMapOvr")
        if not clr_map_nodes:
            raise ValueError("could not find p:clrMapOvr on slide")

        parent = clr_map_nodes[0].getparent()
        transition = etree.Element(
            etree.QName(P_NS, "transition"),
            attrib={"spd": "med"},
        )
        parent.insert(-1, transition)

    transition.set("advTm", str(duration_ms))
    transition.set("advClick", adv_click)
    transition.set("spd", "med")


def timed_output_path(pptx_path: Path) -> Path:
    return pptx_path.with_name(f"{pptx_path.stem}_timed.pptx")


def format_duration(ms: int) -> str:
    seconds = ms / 1000
    return f"{seconds:.1f}s"


def process_slide(
    slide_number: int,
    slide,
    audio_files: Dict[int, Path],
    stats: dict,
) -> None:
    stats["slides_processed"] += 1

    try:
        duration_ms = advance_duration_ms(slide_number, audio_files)
        set_slide_advance(slide, duration_ms)

        if slide_number in audio_files:
            print(
                f"Slide {slide_number}: auto-advance after {format_duration(duration_ms)} "
                f"(MP3 + {BUFFER_MS}ms buffer)"
            )
            stats["slides_with_audio"] += 1
        else:
            print(
                f"Slide {slide_number}: auto-advance after {format_duration(duration_ms)} "
                f"(default, no MP3)"
            )
            stats["slides_without_audio"] += 1

        stats["timings_set"] += 1
    except Exception as exc:
        print(f"Slide {slide_number}: error - {exc}")
        stats["errors"] += 1


def main() -> None:
    OUTPUT_DIR.mkdir(exist_ok=True)
    AUDIO_DIR.mkdir(exist_ok=True)

    pptx_path = find_embedded_powerpoint()
    audio_files = find_audio_files()

    presentation = Presentation(str(pptx_path))
    total_slides = len(presentation.slides)

    if TEST_MODE:
        target_slides = set(TEST_SLIDES)
        print(f"Test mode: setting timings on slides {sorted(target_slides)}")
    else:
        target_slides = None
        print("Full mode: setting timings on all slides")

    stats = {
        "slides_processed": 0,
        "timings_set": 0,
        "slides_with_audio": 0,
        "slides_without_audio": 0,
        "errors": 0,
    }

    for slide_number, slide in enumerate(presentation.slides, start=1):
        if target_slides is not None and slide_number not in target_slides:
            continue
        process_slide(slide_number, slide, audio_files, stats)

    output_path = timed_output_path(pptx_path)
    presentation.save(str(output_path))

    print()
    print(f"Input: {pptx_path.name}")
    print(f"Output: {output_path}")
    print(f"Total slides in deck: {total_slides}")
    print(f"MP3 files available: {len(audio_files)}")
    print(f"Buffer after narration: {BUFFER_MS}ms")
    print(f"Default duration (no MP3): {DEFAULT_SLIDE_DURATION_MS}ms")
    print(f"Click advance allowed: {ALLOW_CLICK_ADVANCE}")

    if TEST_MODE:
        valid_test_slides = [s for s in TEST_SLIDES if 1 <= s <= total_slides]
        print(f"Test slides requested: {TEST_SLIDES}")
        print(f"Test slides in range: {valid_test_slides}")

    print(f"Slides processed: {stats['slides_processed']}")
    print(f"Timings set: {stats['timings_set']}")
    print(f"Slides with MP3 timing: {stats['slides_with_audio']}")
    print(f"Slides with default timing: {stats['slides_without_audio']}")
    print(f"Errors: {stats['errors']}")


if __name__ == "__main__":
    main()
