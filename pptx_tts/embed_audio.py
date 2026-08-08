"""Embed Phase 1 MP3 files into PowerPoint slides (Phase 2).

Auto-plays narration when each slide appears. Advance to the next slide on click.
Does not modify the original .pptx — writes a new file to output/.
"""

import re
import sys
from pathlib import Path
from typing import Dict

from lxml import etree
from pptx import Presentation
from pptx.util import Inches

# ---------------------------------------------------------------------------
# Settings — edit these before running
# ---------------------------------------------------------------------------

AUTO_PLAY = True
HIDE_ICON_DURING_SHOW = True

TEST_MODE = False
TEST_SLIDES = [1, 5, 10]

# Small speaker icon in bottom-right corner
ICON_SIZE = Inches(0.5)
ICON_MARGIN = Inches(0.62)

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

BASE_DIR = Path(__file__).resolve().parent
INPUT_DIR = BASE_DIR / "input"
AUDIO_DIR = BASE_DIR / "audio"
OUTPUT_DIR = BASE_DIR / "output"

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
NSMAP = {"p": P_NS}

AUDIO_FILE_PATTERN = re.compile(r"^slide_(\d+)\.mp3$", re.IGNORECASE)


def find_powerpoint() -> Path:
    """Find the single .pptx file in input/. Exit on 0 or >1 matches."""
    pptx_files = sorted(
        p for p in INPUT_DIR.glob("*.pptx") if not p.name.startswith("~$")
    )

    if not pptx_files:
        print(f"Error: No PowerPoint file found in {INPUT_DIR}")
        print("Place one .pptx file in the input/ folder and run again.")
        sys.exit(1)

    if len(pptx_files) > 1:
        print(f"Error: Found {len(pptx_files)} PowerPoint files in {INPUT_DIR}:")
        for f in pptx_files:
            print(f"  - {f.name}")
        print("Leave only one .pptx file in the input/ folder and run again.")
        sys.exit(1)

    return pptx_files[0]


def find_audio_files() -> Dict[int, Path]:
    """Map slide numbers to MP3 paths from audio/slide_XXX.mp3."""
    audio_files: Dict[int, Path] = {}

    for path in sorted(AUDIO_DIR.glob("slide_*.mp3")):
        match = AUDIO_FILE_PATTERN.match(path.name)
        if match:
            audio_files[int(match.group(1))] = path

    return audio_files


def xpath(element, query: str):
    return etree.ElementBase.xpath(element, query, namespaces=NSMAP)


def shape_element(media_shape):
    return media_shape._element


def set_autoplay(media_shape) -> bool:
    """Configure embedded media to start when the slide appears."""
    element = shape_element(media_shape)
    shape_id = xpath(element, ".//p:cNvPr")[0].attrib["id"]
    slide_element = element.getparent().getparent().getparent()

    for media_tag in ("video", "audio"):
        targets = xpath(
            slide_element,
            f'.//p:timing//p:{media_tag}//p:spTgt[@spid="{shape_id}"]',
        )
        if not targets:
            continue

        cond_nodes = xpath(targets[0].getparent().getparent(), ".//p:cond")
        if cond_nodes:
            cond_nodes[0].set("delay", "0")
            return True

    return False


def hide_media_icon(media_shape) -> bool:
    """Hide the speaker icon during the slideshow."""
    element = shape_element(media_shape)
    shape_id = xpath(element, ".//p:cNvPr")[0].attrib["id"]
    slide_element = element.getparent().getparent().getparent()

    for media_tag in ("video", "audio"):
        nodes = xpath(
            slide_element,
            f'.//p:timing//p:{media_tag}//p:spTgt[@spid="{shape_id}"]/ancestor::p:cMediaNode',
        )
        if nodes:
            nodes[0].set("showWhenStopped", "0")
            return True

    return False


def embed_audio(slide, audio_path: Path, presentation: Presentation):
    """Add MP3 to slide and configure playback."""
    left = presentation.slide_width - ICON_MARGIN - ICON_SIZE
    top = presentation.slide_height - ICON_MARGIN - ICON_SIZE

    media = slide.shapes.add_movie(
        str(audio_path),
        left,
        top,
        ICON_SIZE,
        ICON_SIZE,
        mime_type="audio/mpeg",
    )

    try:
        if AUTO_PLAY and not set_autoplay(media):
            print("  warning: could not enable auto-play for this slide")
    except Exception as exc:
        print(f"  warning: auto-play setup failed ({exc})")

    try:
        if HIDE_ICON_DURING_SHOW and not hide_media_icon(media):
            print("  warning: could not hide icon during show for this slide")
    except Exception as exc:
        print(f"  warning: hide icon setup failed ({exc})")

    return media


def output_path_for(pptx_path: Path) -> Path:
    return OUTPUT_DIR / f"{pptx_path.stem}_with_audio.pptx"


def process_slide(
    slide_number: int,
    slide,
    audio_files: Dict[int, Path],
    presentation: Presentation,
    stats: dict,
) -> None:
    """Embed audio on one slide if an MP3 exists."""
    stats["slides_processed"] += 1

    audio_path = audio_files.get(slide_number)
    if audio_path is None:
        print(f"Slide {slide_number}: no MP3 - skipped")
        stats["slides_without_audio"] += 1
        return

    try:
        embed_audio(slide, audio_path, presentation)
        print(f"Slide {slide_number}: embedded {audio_path.name}")
        stats["audio_embedded"] += 1
    except Exception as exc:
        print(f"Slide {slide_number}: error - {exc}")
        stats["errors"] += 1


def main() -> None:
    INPUT_DIR.mkdir(exist_ok=True)
    AUDIO_DIR.mkdir(exist_ok=True)
    OUTPUT_DIR.mkdir(exist_ok=True)

    pptx_path = find_powerpoint()
    audio_files = find_audio_files()

    if not audio_files:
        print(f"Error: No MP3 files found in {AUDIO_DIR}")
        print("Expected files like slide_001.mp3. Run main.py first.")
        sys.exit(1)

    presentation = Presentation(str(pptx_path))
    total_slides = len(presentation.slides)

    if TEST_MODE:
        target_slides = set(TEST_SLIDES)
        print(f"Test mode: embedding audio on slides {sorted(target_slides)}")
    else:
        target_slides = None
        print("Full mode: embedding audio on all slides with MP3 files")

    stats = {
        "slides_processed": 0,
        "audio_embedded": 0,
        "slides_without_audio": 0,
        "errors": 0,
    }

    for slide_number, slide in enumerate(presentation.slides, start=1):
        if target_slides is not None and slide_number not in target_slides:
            continue
        process_slide(slide_number, slide, audio_files, presentation, stats)

    output_path = output_path_for(pptx_path)
    presentation.save(str(output_path))

    print()
    print(f"PowerPoint: {pptx_path.name}")
    print(f"Output: {output_path}")
    print(f"Total slides in deck: {total_slides}")
    print(f"MP3 files available: {len(audio_files)}")

    if TEST_MODE:
        valid_test_slides = [s for s in TEST_SLIDES if 1 <= s <= total_slides]
        print(f"Test slides requested: {TEST_SLIDES}")
        print(f"Test slides in range: {valid_test_slides}")

    print(f"Slides processed: {stats['slides_processed']}")
    print(f"Audio embedded: {stats['audio_embedded']}")
    print(f"Slides without MP3: {stats['slides_without_audio']}")
    print(f"Errors: {stats['errors']}")


if __name__ == "__main__":
    main()
